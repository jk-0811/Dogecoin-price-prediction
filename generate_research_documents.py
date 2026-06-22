"""Generate research-style PDF, DOCX, and PPT deliverables for DOGE prediction."""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import List

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.util import Inches as PPTInches
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import Image as RLImage
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
from sklearn.ensemble import RandomForestRegressor
from sklearn.metrics import mean_absolute_error, mean_squared_error, r2_score
from sklearn.model_selection import train_test_split
from sklearn.preprocessing import StandardScaler


@dataclass
class TrainingResult:
    raw_rows: int
    processed_rows: int
    start_date: str
    end_date: str
    train_size: int
    test_size: int
    mse: float
    rmse: float
    mae: float
    r2: float
    features: List[str]
    importances: List[tuple[str, float]]
    sample_outputs: List[tuple[float, float]]
    close_series: pd.DataFrame


def train_model(csv_path: Path) -> TrainingResult:
    df = pd.read_csv(csv_path)
    raw_rows = len(df)

    df["Date"] = pd.to_datetime(df["Date"])
    df = df.sort_values("Date").dropna().reset_index(drop=True)
    df["MA5"] = df["Close"].rolling(window=5).mean()
    df["MA10"] = df["Close"].rolling(window=10).mean()
    df["Daily_Return"] = df["Close"].pct_change()
    df["Volatility"] = df["Close"].rolling(window=10).std()
    df = df.dropna().reset_index(drop=True)

    feature_cols = ["Open", "High", "Low", "Volume", "MA5", "MA10", "Daily_Return", "Volatility"]
    x = df[feature_cols]
    y = df["Close"]

    scaler = StandardScaler()
    x_scaled = scaler.fit_transform(x)
    x_train, x_test, y_train, y_test = train_test_split(
        x_scaled, y, test_size=0.2, random_state=42
    )

    model = RandomForestRegressor(
        n_estimators=100,
        max_depth=15,
        min_samples_split=5,
        min_samples_leaf=2,
        random_state=42,
        n_jobs=-1,
    )
    model.fit(x_train, y_train)
    y_pred = model.predict(x_test)

    mse = mean_squared_error(y_test, y_pred)
    rmse = float(np.sqrt(mse))
    mae = mean_absolute_error(y_test, y_pred)
    r2 = r2_score(y_test, y_pred)

    importances = sorted(
        zip(feature_cols, model.feature_importances_), key=lambda item: item[1], reverse=True
    )

    sample_outputs = [
        (float(actual), float(predicted))
        for actual, predicted in list(zip(y_test.iloc[:5], y_pred[:5]))
    ]

    close_series = df[["Date", "Close"]].copy()

    return TrainingResult(
        raw_rows=raw_rows,
        processed_rows=len(df),
        start_date=str(df["Date"].min().date()),
        end_date=str(df["Date"].max().date()),
        train_size=len(x_train),
        test_size=len(x_test),
        mse=float(mse),
        rmse=rmse,
        mae=float(mae),
        r2=float(r2),
        features=feature_cols,
        importances=[(name, float(score)) for name, score in importances],
        sample_outputs=sample_outputs,
        close_series=close_series,
    )


def generate_charts(result: TrainingResult, output_dir: Path) -> tuple[Path, Path]:
    trend_path = output_dir / "dataset_trend.png"
    plt.figure(figsize=(10, 4))
    plt.plot(result.close_series["Date"], result.close_series["Close"], color="#1f77b4", linewidth=1.2)
    plt.title("DOGE Close Price Trend")
    plt.xlabel("Date")
    plt.ylabel("Close Price (USD)")
    plt.tight_layout()
    plt.savefig(trend_path, dpi=200)
    plt.close()

    scatter_path = output_dir / "actual_vs_predicted.png"
    actual_values = [pair[0] for pair in result.sample_outputs]
    predicted_values = [pair[1] for pair in result.sample_outputs]
    plt.figure(figsize=(6, 4))
    plt.scatter(actual_values, predicted_values, color="#2ca02c", s=55)
    min_val = min(actual_values + predicted_values)
    max_val = max(actual_values + predicted_values)
    plt.plot([min_val, max_val], [min_val, max_val], color="#d62728", linestyle="--")
    plt.title("Sample Actual vs Predicted")
    plt.xlabel("Actual Close")
    plt.ylabel("Predicted Close")
    plt.tight_layout()
    plt.savefig(scatter_path, dpi=200)
    plt.close()

    return trend_path, scatter_path


def build_docx(result: TrainingResult, trend_img: Path, scatter_img: Path, file_path: Path) -> None:
    doc = Document()

    title = doc.add_heading("Dogecoin Price Prediction using Random Forest Regression", 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle = doc.add_paragraph("Research Summary Report")
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle.runs[0].italic = True

    doc.add_heading("Abstract", level=1)
    doc.add_paragraph(
        "This report presents a supervised machine learning framework for predicting DOGE-USD close prices "
        "using engineered technical indicators and market variables. A Random Forest Regressor was trained "
        "on historical data and evaluated using holdout validation metrics."
    )

    doc.add_heading("1. Problem Definition", level=1)
    doc.add_paragraph(
        "The objective is to estimate the next close price of Dogecoin using recent market context. "
        "The model takes open, high, low, volume, and derived indicators as input features."
    )

    doc.add_heading("2. Dataset Used", level=1)
    doc.add_paragraph(
        f"Source file: DOGE-USD.csv | Raw rows: {result.raw_rows} | Processed rows: {result.processed_rows}"
    )
    doc.add_paragraph(
        f"Date range: {result.start_date} to {result.end_date} | Train/Test split: {result.train_size}/{result.test_size}"
    )
    doc.add_picture(str(trend_img), width=Inches(6.2))

    doc.add_heading("3. Model Trained", level=1)
    doc.add_paragraph("Algorithm: RandomForestRegressor")
    doc.add_paragraph(
        "Hyperparameters: n_estimators=100, max_depth=15, min_samples_split=5, min_samples_leaf=2, random_state=42"
    )
    doc.add_paragraph("Input Features: " + ", ".join(result.features))

    doc.add_heading("4. Accuracy and Evaluation", level=1)
    metrics_table = doc.add_table(rows=1, cols=2)
    metrics_table.style = "Light Grid Accent 1"
    hdr = metrics_table.rows[0].cells
    hdr[0].text = "Metric"
    hdr[1].text = "Value"
    for metric_name, metric_value in [
        ("MSE", f"{result.mse:.8f}"),
        ("RMSE", f"{result.rmse:.8f}"),
        ("MAE", f"{result.mae:.8f}"),
        ("R2", f"{result.r2:.6f}"),
    ]:
        row = metrics_table.add_row().cells
        row[0].text = metric_name
        row[1].text = metric_value

    doc.add_heading("5. Output Visibility", level=1)
    out_table = doc.add_table(rows=1, cols=3)
    out_table.style = "Light Grid Accent 2"
    out_hdr = out_table.rows[0].cells
    out_hdr[0].text = "Sample"
    out_hdr[1].text = "Actual Close"
    out_hdr[2].text = "Predicted Close"
    for idx, (actual, predicted) in enumerate(result.sample_outputs, start=1):
        row = out_table.add_row().cells
        row[0].text = str(idx)
        row[1].text = f"{actual:.6f}"
        row[2].text = f"{predicted:.6f}"

    doc.add_picture(str(scatter_img), width=Inches(5.7))
    caption = doc.add_paragraph("Figure: Sample Actual vs Predicted values")
    caption.alignment = WD_ALIGN_PARAGRAPH.CENTER

    for paragraph in doc.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Times New Roman"
            run.font.size = Pt(11)

    doc.save(file_path)


def build_ppt(result: TrainingResult, trend_img: Path, scatter_img: Path, file_path: Path) -> None:
    presentation = Presentation()

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = "Dogecoin Price Prediction"
    title_slide.placeholders[1].text = "Research Presentation | Random Forest Regression"

    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Problem Definition"
    slide.placeholders[1].text = (
        "Goal: Predict DOGE close price from market and technical indicators.\n"
        "Inputs: Open, High, Low, Volume, MA5, MA10, Daily Return, Volatility.\n"
        "Task type: Supervised regression."
    )

    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Dataset Used"
    text_box = slide.shapes.add_textbox(PPTInches(0.5), PPTInches(0.9), PPTInches(4.2), PPTInches(2.0))
    frame = text_box.text_frame
    frame.text = f"Rows (raw): {result.raw_rows}"
    frame.add_paragraph().text = f"Rows (processed): {result.processed_rows}"
    frame.add_paragraph().text = f"Date range: {result.start_date} to {result.end_date}"
    frame.add_paragraph().text = f"Split: {result.train_size} train / {result.test_size} test"
    slide.shapes.add_picture(str(trend_img), PPTInches(4.8), PPTInches(0.9), width=PPTInches(4.3))

    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Model Trained"
    slide.placeholders[1].text = (
        "RandomForestRegressor\n"
        "n_estimators=100, max_depth=15\n"
        "min_samples_split=5, min_samples_leaf=2\n"
        "random_state=42"
    )

    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Accuracy Shown"
    metrics_box = slide.shapes.add_textbox(PPTInches(0.5), PPTInches(1.0), PPTInches(4.0), PPTInches(2.3))
    metrics_frame = metrics_box.text_frame
    metrics_frame.text = f"MSE: {result.mse:.8f}"
    metrics_frame.add_paragraph().text = f"RMSE: {result.rmse:.8f}"
    metrics_frame.add_paragraph().text = f"MAE: {result.mae:.8f}"
    metrics_frame.add_paragraph().text = f"R2: {result.r2:.6f}"
    slide.shapes.add_picture(str(scatter_img), PPTInches(4.8), PPTInches(1.0), width=PPTInches(4.3))

    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Output Visible"
    output_box = slide.shapes.add_textbox(PPTInches(0.5), PPTInches(1.0), PPTInches(8.8), PPTInches(4.5))
    output_frame = output_box.text_frame
    output_frame.text = "Sample predictions from test data:"
    for idx, (actual, predicted) in enumerate(result.sample_outputs, start=1):
        output_frame.add_paragraph().text = (
            f"Sample {idx}: Actual={actual:.6f} | Predicted={predicted:.6f}"
        )

    presentation.save(file_path)


def build_pdf(result: TrainingResult, trend_img: Path, scatter_img: Path, file_path: Path) -> None:
    doc = SimpleDocTemplate(str(file_path), pagesize=A4, rightMargin=2 * cm, leftMargin=2 * cm)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "TitleStyle",
        parent=styles["Title"],
        fontName="Times-Bold",
        fontSize=18,
        spaceAfter=14,
    )
    heading_style = ParagraphStyle(
        "HeadingStyle",
        parent=styles["Heading2"],
        fontName="Times-Bold",
        fontSize=13,
        spaceBefore=10,
        spaceAfter=6,
    )
    body_style = ParagraphStyle(
        "BodyStyle", parent=styles["BodyText"], fontName="Times-Roman", fontSize=10, leading=14
    )

    story = []
    story.append(Paragraph("Dogecoin Price Prediction using Random Forest Regression", title_style))
    story.append(Paragraph("Research Paper Summary", body_style))
    story.append(Spacer(1, 10))

    story.append(Paragraph("Abstract", heading_style))
    story.append(
        Paragraph(
            "This paper investigates a machine learning approach for forecasting DOGE-USD close prices. "
            "The framework combines market inputs with technical indicators and evaluates a Random Forest "
            "Regressor using holdout testing metrics.",
            body_style,
        )
    )

    story.append(Paragraph("1. Problem Defined", heading_style))
    story.append(
        Paragraph(
            "Predict the close price of Dogecoin from open/high/low/volume values and engineered signals "
            "(MA5, MA10, Daily Return, Volatility).",
            body_style,
        )
    )

    story.append(Paragraph("2. Dataset Used", heading_style))
    story.append(
        Paragraph(
            f"DOGE-USD.csv with {result.raw_rows} raw rows and {result.processed_rows} processed rows; "
            f"date window {result.start_date} to {result.end_date}; split {result.train_size}/{result.test_size}.",
            body_style,
        )
    )
    story.append(Spacer(1, 6))
    story.append(RLImage(str(trend_img), width=15 * cm, height=6 * cm))

    story.append(Paragraph("3. Model Trained", heading_style))
    story.append(
        Paragraph(
            "RandomForestRegressor configured with 100 estimators, max depth 15, min split 5, min leaf 2.",
            body_style,
        )
    )

    story.append(Paragraph("4. Accuracy Shown", heading_style))
    metric_data = [
        ["Metric", "Value"],
        ["MSE", f"{result.mse:.8f}"],
        ["RMSE", f"{result.rmse:.8f}"],
        ["MAE", f"{result.mae:.8f}"],
        ["R2", f"{result.r2:.6f}"],
    ]
    metrics_table = Table(metric_data, colWidths=[6 * cm, 6 * cm])
    metrics_table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1f4e78")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("GRID", (0, 0), (-1, -1), 0.6, colors.grey),
                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                ("FONTNAME", (0, 0), (-1, 0), "Times-Bold"),
            ]
        )
    )
    story.append(metrics_table)

    story.append(Paragraph("5. Output Visible", heading_style))
    output_data = [["Sample", "Actual", "Predicted"]]
    for idx, (actual, predicted) in enumerate(result.sample_outputs, start=1):
        output_data.append([str(idx), f"{actual:.6f}", f"{predicted:.6f}"])
    output_table = Table(output_data, colWidths=[3 * cm, 4 * cm, 4 * cm])
    output_table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2e8b57")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("GRID", (0, 0), (-1, -1), 0.6, colors.grey),
                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                ("FONTNAME", (0, 0), (-1, 0), "Times-Bold"),
            ]
        )
    )
    story.append(output_table)
    story.append(Spacer(1, 6))
    story.append(RLImage(str(scatter_img), width=10 * cm, height=6 * cm))

    doc.build(story)


def main() -> None:
    base_dir = Path(__file__).resolve().parent
    csv_path = base_dir / "DOGE-USD.csv"
    output_dir = base_dir / "research_outputs"
    output_dir.mkdir(parents=True, exist_ok=True)

    result = train_model(csv_path)
    trend_img, scatter_img = generate_charts(result, output_dir)

    pdf_path = output_dir / "Dogecoin_Research_Paper.pdf"
    docx_path = output_dir / "Dogecoin_Research_Paper.docx"
    ppt_path = output_dir / "Dogecoin_Research_Presentation.pptx"

    build_pdf(result, trend_img, scatter_img, pdf_path)
    build_docx(result, trend_img, scatter_img, docx_path)
    build_ppt(result, trend_img, scatter_img, ppt_path)

    print("Generated files:")
    print(f"- {pdf_path}")
    print(f"- {docx_path}")
    print(f"- {ppt_path}")
    print("\nKey metrics:")
    print(f"MSE={result.mse:.8f}")
    print(f"RMSE={result.rmse:.8f}")
    print(f"MAE={result.mae:.8f}")
    print(f"R2={result.r2:.6f}")


if __name__ == "__main__":
    main()
